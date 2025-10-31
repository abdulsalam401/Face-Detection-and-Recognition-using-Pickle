# import tensorflow as tf

# if tf.config.list_physical_devices('GPU'):
#     print("GPU is available!")
# else:
#     print("GPU is not available.")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Helper function to add slide with title and bullet points
def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = '\n'.join(bullets)

# Create presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Business Model: Compressed Sofa Import & Sales"
slide.placeholders[1].text = "By Abdul Salam"

# Business Model Canvas slides
add_bullet_slide(prs, "Key Partners", [
    "Furniture manufacturers in China",
    "Shipping/logistics companies",
    "Local customs agents",
    "Retail partners (showrooms or dealers)"
])

add_bullet_slide(prs, "Key Activities", [
    "Importing and shipping compressed sofas",
    "Storage and warehousing",
    "Marketing and advertising locally",
    "Direct and wholesale selling",
    "Customer service and returns"
])

add_bullet_slide(prs, "Key Resources", [
    "Capital for importing and logistics",
    "Storage facility",
    "Website or online store",
    "Sales team or network",
    "Transportation for delivery"
])

add_bullet_slide(prs, "Value Propositions", [
    "Affordable, stylish furniture",
    "Space-saving design (compressed sofas)",
    "Modern and trendy looks",
    "Easy to transport and assemble",
    "Lower price compared to local furniture"
])

add_bullet_slide(prs, "Customer Relationships", [
    "After-sale support",
    "Social media interaction",
    "Discounts and loyalty programs",
    "Installation support"
])

add_bullet_slide(prs, "Channels", [
    "Physical store/showroom",
    "Online store (website, Facebook, Daraz)",
    "WhatsApp and phone orders",
    "Local ads and exhibitions"
])

add_bullet_slide(prs, "Customer Segments", [
    "Newly married couples",
    "University students",
    "Middle-income families",
    "Renters and small apartment dwellers"
])

add_bullet_slide(prs, "Cost Structure (PKR)", [
    "Product Cost (per unit): PKR 18,000",
    "Shipping & Customs: PKR 6,000",
    "Storage Rent (monthly): PKR 25,000",
    "Marketing (monthly): PKR 15,000",
    "Local Delivery (per unit): PKR 2,000"
])

add_bullet_slide(prs, "Revenue Streams (PKR)", [
    "Retail Price (per unit): PKR 32,000",
    "Estimated Monthly Sales: 50 units",
    "Monthly Revenue: PKR 1,600,000",
    "Monthly Cost: PKR 1,225,000",
    "Monthly Profit: PKR 375,000"
])

# Save the presentation
pptx_path = "/mnt/data/Compressed_Sofa_Business_Model.pptx"
prs.save(pptx_path)

pptx_path
